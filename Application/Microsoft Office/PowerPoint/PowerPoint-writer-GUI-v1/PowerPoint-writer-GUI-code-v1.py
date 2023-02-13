import PySimpleGUI as sg
import pptx 
from pptx.util import Inches,Pt
from pptx.enum.text import PP_ALIGN,PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import PP_PLACEHOLDER,PP_PLACEHOLDER_TYPE

import os

class PPT_Window():
    
    def MainWindow(self):
        
        def ClearAllValues():
            for k,v in values.items():
                window[k].update(emptyString)
        
        def UpdateData():
            
            ##
            # Set the list (NOT use hardcode)
            List1=["Title","Content"]
            List2=["Slide","Page"]
            
            # Combine all possibles of List1 and List2 and store it into List3
            List3=list()
            for x in List2:
               for y in List1:  
                  List3.append(x+" of "+y+":")
           
            List4=["Bold?","Underline?","Italic?","Strike?","Double Strike?"]
            List5=["Input Value:"]
            
            otherList=["Document name( NOT including file extension):","Document Location:","Number of pages(NOT including first slide):"]
            ##
            
            updateMessage=""
            window[message_key].update(updateMessage)
            
            ## Counter
            # Count elem for List4
            innerLoopCnt=0
            # Count elem for List3
            outerLoopCnt=0
            # Counter to count the rest values
            restCnt=0
            
            ##
            
            ## Variable to record the status of skip.
            # Current counter for skip.
            skipCurCnt=0
            # Total time we have to skip.
            skipTotalCnt=1
            # Determine the next time we have to skip.
            startToSkip=False
            # Has the skip already started?
            hasSkipped=False
            ##
            
            # Update all values in the window.
            for k,v in values.items():
                
                # The settings of font settings and input values for input fields for each row.
                # They are 4 rows. They are for List3 respectively.
                if outerLoopCnt<len(List3):         
                    
                    # Check it is the beginning for each row.
                    # If the condition are true, then it is the beginning for each row.
                    if outerLoopCnt<len(List3) and innerLoopCnt==0 and hasSkipped==False:    
                        updateMessage=updateMessage+List3[outerLoopCnt]
                        updateMessage=updateMessage+"\n"
                        startToSkip=True
                        hasSkipped=True      
                        
                    # Check to have to start to skip
                    if startToSkip==True:
                        # Skip the checkboxes and input value of input fields.
                        updateMessage=updateMessage+List5[skipCurCnt]
                        updateMessage=updateMessage+str(v)
                        updateMessage=updateMessage+"\n"
                        skipCurCnt=skipCurCnt+1
                        
                        # Check we have to continue to skip, if the conditional is true, we don't have to.
                        if skipCurCnt>=skipTotalCnt:
                            skipCurCnt=0
                            startToSkip=False
                            hasSkipped=True
                    else:
                        # Update all values for the checkboxes (they are about font settings).
                        updateMessage=updateMessage+List4[innerLoopCnt]   
                        updateMessage=updateMessage+str(v)
                        updateMessage=updateMessage+"\t"
                        innerLoopCnt=innerLoopCnt+1
                        
                        # Check the all possibility for checkboxes are scanned.
                        if innerLoopCnt== len(List4):     
                            innerLoopCnt=0
                            outerLoopCnt=outerLoopCnt+1
                            hasSkipped=False
                            updateMessage=updateMessage+"\n"
                    
                # Rest
                else:
                    updateMessage=updateMessage+"\n"
                    updateMessage=updateMessage+otherList[restCnt]
                    restCnt=restCnt+1
                    updateMessage=updateMessage+str(v)
                    
                    
            updateMessage=updateMessage+"\n"
            window[message_key].update(updateMessage)
         
        def PopupErrorWindow(title:str,e:str):
            sg.popup_error(e,title=title,non_blocking=False)
            
        def PopupSaveFileWindow():
            emptyString=""
            while True:
                fileLocation=sg.popup_get_folder(message="Find a location for saving auto-generated files.",title="Location Select Window")
                
                
                if fileLocation is None or fileLocation == emptyString:
                    break
                
                ok=True
                if ok and not os.path.exists(fileLocation):
                    ok=False
                    title="Error Window"
                    e="ERROR!!! File Location does NOT exists!!!"
                    PopupErrorWindow(e=e,title=title)
                if ok and not os.path.isdir(fileLocation):
                    ok=False
                    title="Error Window"
                    e="ERROR!!! File Location is NOT a directory!!!"
                    PopupErrorWindow(e=e,title=title)
                
                if ok:
                    return fileLocation
                
            return emptyString
        
        
            
        def CreatePPT():
            
            FetchData()
            
            slideTitleMessage=values[slideTitle_key]
            slideContentMessage=values[slideContent_key]
            
            pageTitleMessage=values[pageTitle_key]
            pageContentMessage=values[pageContent_key]
            
            mypptx=pptx.Presentation()
            
            slide_layout=mypptx.slide_layouts[0]
            mypptx_slide=mypptx.slides.add_slide(slide_layout)

            mypptx_slide.placeholders[0].text=slideTitleMessage
     
            mypptx_slide.placeholders[0].text_frame.paragraphs[0].font.bold=slideTitleFontGroup.bold
            mypptx_slide.placeholders[0].text_frame.paragraphs[0].font.underline=slideTitleFontGroup.underline
            mypptx_slide.placeholders[0].text_frame.paragraphs[0].font.italic=slideTitleFontGroup.italic
            mypptx_slide.placeholders[0].text_frame.paragraphs[0].font.strike=slideTitleFontGroup.strike
            mypptx_slide.placeholders[0].text_frame.paragraphs[0].font.double_strike=slideTitleFontGroup.double_strike
            
            mypptx_slide.placeholders[1].text=slideContentMessage
            
            mypptx_slide.placeholders[1].text_frame.paragraphs[0].font.bold=slideContentFontGroup.bold
            mypptx_slide.placeholders[1].text_frame.paragraphs[0].font.underline=slideContentFontGroup.underline
            mypptx_slide.placeholders[1].text_frame.paragraphs[0].font.italic=slideContentFontGroup.italic
            mypptx_slide.placeholders[1].text_frame.paragraphs[0].font.strike=slideContentFontGroup.strike
            mypptx_slide.placeholders[1].text_frame.paragraphs[0].font.double_strike=slideContentFontGroup.double_strike
        
            for i in range(0,int(values[numOfPage_key],10),1):
                slide_layout=mypptx.slide_layouts[5]
                mypptx_slide=mypptx.slides.add_slide(slide_layout)
                mypptx_slide.placeholders[0].text=pageTitleMessage
                mypptx_slide.placeholders[0].text_frame.paragraphs[0].font.bold=pageTitleFontGroup.bold
                mypptx_slide.placeholders[0].text_frame.paragraphs[0].font.underline=pageTitleFontGroup.underline
                mypptx_slide.placeholders[0].text_frame.paragraphs[0].font.italic=pageTitleFontGroup.italic
                mypptx_slide.placeholders[0].text_frame.paragraphs[0].font.strike=pageTitleFontGroup.strike
                mypptx_slide.placeholders[0].text_frame.paragraphs[0].font.double_strike=pageTitleFontGroup.double_strike
            
                # For adjusting the  Margins in inches 
                left = Inches(1)
                width= Inches(1)
                top = Inches(2)
                height = Inches(1)
                
                # creating textBox
                txBox = mypptx_slide.shapes.add_textbox(left, top,width, height)
  
                # creating textFrames
                tf = txBox.text_frame
  
                # adding Paragraphs
                p = tf.add_paragraph()
                p.font.bold = pageContentFontGroup.bold
                p.font.underline= pageContentFontGroup.underline
                p.font.italic = pageContentFontGroup.italic
                p.font.strike = pageContentFontGroup.strike
                p.font.double_strike = pageContentFontGroup.double_strike
  
                # adding text
                p.text = pageContentMessage
            
            
            
            filename=values[docName_key]
            filepath=values[docLoc_key]
            fileext=".pptx"
            filepathname=os.path.join(filepath,filename+str(fileext))
            mypptx.save(filepathname)
            
        ##
        # Some definitions about strings.
        emptyString=""
        #Some keys definition of input fields.
        
        slideTitle_key="--Slide TITLE--"
        slideContent_key="--Slide CONTENT--"
        pageTitle_key="--Page TITLE--"
        pageContent_key="--Page CONTENT--"
        
        docName_key="--Document Name--"
        docLoc_key="--Document Location--"
        numOfPage_key="--Number of Pages--"
        
        message_key="-Message-"
        
        #Some keys definition of checkboxes.
        
        class FontGroup():
            def __init__(self,name):
                self.name = name
                self.bold=True
                self.underline=True
                self.italic=True
                self.strike=True
                self.double_strike=True
                
                self.BoldName="Bold?"
                self.UnderlineName="Underline?"
                self.ItalicName="Italic?"
                self.StrikeName="Strike?"
                self.DoubleStrikeName="Double Strike?"
                
        
            def __dict__(self):
                dic={
                "self.bold":self.bold
                ,"self.underline":self.underline
                ,"self.italic":self.italic
                ,"self.strike":self.strike
                ,"self.double_strike":self.double_strike
                ,"self.BoldName":self.BoldName
                ,"self.UnderlineName":self.UnderlineName
                ,"self.ItalicName":self.ItalicName
                ,"self.StrikeName":self.StrikeName
                ,"self.DoubleStrikeName":self.DoubleStrikeName
                ,"-BOLD-":("Bold "+self.name)
                ,"-UNDERLINE-":("Underline "+self.name)
                ,"-STRIKE-":("Stike "+self.name)
                ,"-ITALIC-":("Italic "+self.name)
                ,"-DOUBLE_STRIKE-":("Double Strike "+self.name)
                }
                return dic
            
            def UpdateFontSettings(self):
                self.bold=values[self.__dict__()[BOLD_key]]
                self.underline=values[self.__dict__()[UNDERLINE_key]]
                self.italic=values[self.__dict__()[ITALIC_key]]
                self.strike=values[self.__dict__()[STRIKE_key]]
                self.double_strike=values[self.__dict__()[DOUBLE_STRIKE_key]]
            def Get(self,name):
                return self.__dict__()[name]
                
                
        def FetchData():
            slideTitleFontGroup.UpdateFontSettings()
            slideContentFontGroup.UpdateFontSettings()
            pageTitleFontGroup.UpdateFontSettings()
            pageContentFontGroup.UpdateFontSettings()
            
        slideTitleFontGroup=FontGroup("Slide Title Font")  
        slideContentFontGroup=FontGroup("Slide Content Font")  
        pageTitleFontGroup=FontGroup("Page Title Font")  
        pageContentFontGroup=FontGroup("Page Content Font")  
    
        BOLD_text="BOLD?"
        UNDERLINE_text="UNDERLINE?"
        STRIKE_text="STRIKE?"
        ITALIC_text="ITALIC?"
        DOUBLE_STRIKE_text="DOUBLE STRIKE?"
        
        BOLD_key="-BOLD-"
        UNDERLINE_key="-UNDERLINE-"
        STRIKE_key="-STRIKE-"
        ITALIC_key="-ITALIC-"
        DOUBLE_STRIKE_key="-DOUBLE_STRIKE-"
        
        
        
        #Some keys definition of buttons.
        
        clearAllButton="Clear All"
        createButton="Create"
        closeButton="Close"
          
        browseDocLocButton="Browse Document Location"
        browseDocLocButton_text="Browse"
        ##
        
        slideTitleMessage=emptyString
        slideContentMessage=emptyString
        
        layout=[
            [
                sg.Text("Title of first slide:"),
                sg.Input(key=slideTitle_key),
                sg.Checkbox(text=BOLD_text,key=slideTitleFontGroup.__dict__()[BOLD_key]),
                sg.Checkbox(text=UNDERLINE_text,key=slideTitleFontGroup.__dict__()[UNDERLINE_key]),
                sg.Checkbox(text=STRIKE_text,key=slideTitleFontGroup.__dict__()[STRIKE_key]),
                sg.Checkbox(text=ITALIC_text,key=slideTitleFontGroup.__dict__()[ITALIC_key]),
                sg.Checkbox(text=DOUBLE_STRIKE_text,key=slideTitleFontGroup.__dict__()[DOUBLE_STRIKE_key])
            ],
            [
                sg.Text("Content of the first slide:"),
                sg.Input(key=slideContent_key),
                sg.Checkbox(text=BOLD_text,key=slideContentFontGroup.__dict__()[BOLD_key]),
                sg.Checkbox(text=UNDERLINE_text,key=slideContentFontGroup.__dict__()[UNDERLINE_key]),
                sg.Checkbox(text=STRIKE_text,key=slideContentFontGroup.__dict__()[STRIKE_key]),
                sg.Checkbox(text=ITALIC_text,key=slideContentFontGroup.__dict__()[ITALIC_key]),
                sg.Checkbox(text=DOUBLE_STRIKE_text,key=slideContentFontGroup.__dict__()[DOUBLE_STRIKE_key])
            ],
            [
                sg.Text("Title of the page slide:"),
                sg.Input(key=pageTitle_key),
                sg.Checkbox(text=BOLD_text,key=pageTitleFontGroup.__dict__()[BOLD_key]),
                sg.Checkbox(text=UNDERLINE_text,key=pageTitleFontGroup.__dict__()[UNDERLINE_key]),
                sg.Checkbox(text=STRIKE_text,key=pageTitleFontGroup.__dict__()[STRIKE_key]),
                sg.Checkbox(text=ITALIC_text,key=pageTitleFontGroup.__dict__()[ITALIC_key]),
                sg.Checkbox(text=DOUBLE_STRIKE_text,key=pageTitleFontGroup.__dict__()[DOUBLE_STRIKE_key])
            ],
            [
                sg.Text("Content of the page slide:"),
                sg.Input(key=pageContent_key),
                sg.Checkbox(text=BOLD_text,key=pageContentFontGroup.__dict__()[BOLD_key]),
                sg.Checkbox(text=UNDERLINE_text,key=pageContentFontGroup.__dict__()[UNDERLINE_key]),
                sg.Checkbox(text=STRIKE_text,key=pageContentFontGroup.__dict__()[STRIKE_key]),
                sg.Checkbox(text=ITALIC_text,key=pageContentFontGroup.__dict__()[ITALIC_key]),
                sg.Checkbox(text=DOUBLE_STRIKE_text,key=pageContentFontGroup.__dict__()[DOUBLE_STRIKE_key])
            ],
            [
                sg.Text("Document name:"),
                sg.Input(key=docName_key)
            ],
            [
                 sg.Text("Document location:"),
                 sg.Input(key=docLoc_key,readonly=True),
                 sg.Button(browseDocLocButton_text,key=browseDocLocButton,size=(20,1))
            ],
            [
                 sg.Text("Number of pages(NOT including the first slide):"),
                 sg.Input(key=numOfPage_key)
            ],
            [
                sg.Button(createButton),
                sg.Button(clearAllButton),
                sg.Button(closeButton)    
            ]
            ,
            [
                sg.HSeparator()
            ]
            ,
            [
                 sg.Text("",key=message_key) 
            ]
        ]
        window=sg.Window("PPT writer Window",layout=layout)
        
        
        first=True
        while True:
            event,values=window.read(timeout=1000)
            
            if event == sg.WIN_CLOSED or event == closeButton:
                print(values)
                break
            
            if event == browseDocLocButton:
                fileLocation=PopupSaveFileWindow()
                if fileLocation!=emptyString:
                    window[docLoc_key].update(fileLocation)
                
            if event == clearAllButton:
                ClearAllValues()
                
            if event == createButton:
                CreatePPT()
            
            UpdateData()
            
        window.close()
            
        

def main():
    inst=PPT_Window()
    inst.MainWindow()
    
main()
